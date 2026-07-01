#!/usr/bin/env python3
"""Validate PDF and PPTX derivatives for ms-presentation-deck outputs.

The HTML deck is the source of truth. This gate confirms that derivatives are
stored in the deck derivative folders, match the slide count, preserve the
requested locale's speaker notes, and keep PPTX slides editable.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

LOCALE_MARKERS = {
    "en": ("OPENING", "CORE", "TRANSITION", "TIMING"),
    "pt-BR": ("ABERTURA", "NÚCLEO", "TRANSIÇÃO", "TIMING"),
    "es": ("APERTURA", "NÚCLEO", "TRANSICIÓN", "TIMING"),
}

PLACEHOLDERS = (
    "This slide needs a complete speaker note",
    "Use this slide as a complete presenter prompt",
    "Nota base do slide",
    "Nota base de la diapositiva",
    "[source]",
    "[citation needed]",
)


def fail(message: str) -> None:
    print(f"FAIL: {message}")
    raise SystemExit(1)


def load_i18n(html: str) -> dict:
    marker = "const I18N = "
    if marker not in html:
        fail("HTML source does not contain `const I18N = ...`")
    start = html.index(marker) + len(marker)
    try:
        data, _ = json.JSONDecoder().raw_decode(html[start:])
    except json.JSONDecodeError as exc:
        fail(f"could not parse I18N JSON: {exc}")
    return data


def slide_count(html: str) -> int:
    count = len(re.findall(r"<section\b(?=[^>]*\bclass=[\"'][^\"']*\bslide\b)", html, flags=re.I))
    if count == 0:
        fail("HTML source has no slide sections")
    return count


def visible_i18n_keys(html: str) -> list[str]:
    sections = re.findall(
        r"<section\b(?=[^>]*\bclass=[\"'][^\"']*\bslide\b).*?</section>",
        html,
        flags=re.I | re.S,
    )
    keys: list[str] = []
    for section in sections:
        keys.extend(re.findall(r"data-i18n=[\"']([^\"']+)[\"']", section))
        keys.extend(re.findall(r"data-i18n-list=[\"']([^\"']+)[\"']", section))
    return sorted(set(keys))


def resolve_key(data: dict, locale: str, key: str):
    current = data.get(locale, {})
    for part in key.split("."):
        current = current.get(part, "") if isinstance(current, dict) else ""
    return current


def validate_visible_locale(html: str, locale: str) -> None:
    data = load_i18n(html)
    keys = visible_i18n_keys(html)
    if locale != "en" and not keys:
        fail(f"HTML has no visible data-i18n keys, cannot export localized {locale} derivative")
    missing = [key for key in keys if resolve_key(data, "en", key) not in (None, "") and resolve_key(data, locale, key) in (None, "")]
    if missing:
        fail(f"visible {locale} translations missing for keys: {missing[:12]}")


def note_keys(html: str, total: int) -> list[str]:
    match = re.search(r"const NOTE_MAP = (\[[^\]]*\]);", html, flags=re.S)
    if not match:
        return [f"s{i}" for i in range(1, total + 1)]
    note_map = json.loads(match.group(1))
    return [note_map[i - 1] or f"s{i}" if i - 1 < len(note_map) else f"s{i}" for i in range(1, total + 1)]


def normalize(text: str) -> str:
    text = (text or "").replace("**", "").replace("*", "")
    return re.sub(r"\s+", " ", text).strip()


def strip_markers(text: str) -> str:
    return normalize(re.sub(r"\[[^\]]+\]", "", text))


def expected_notes(html: str, locale: str, total: int) -> list[str]:
    data = load_i18n(html)
    if locale not in data:
        fail(f"locale {locale} missing from I18N")
    notes = data[locale].get("notes", {})
    keys = note_keys(html, total)
    expected = []
    for index, key in enumerate(keys, start=1):
        note = notes.get(key, "")
        if not note:
            fail(f"missing {locale} speaker note {key} for slide {index}")
        if any(placeholder in note for placeholder in PLACEHOLDERS):
            fail(f"placeholder speaker note remains in {locale} note {key}")
        expected.append(note)
    extras = sorted(set(notes) - set(keys))
    if extras:
        fail(f"extra unused note keys for {locale}: {extras[:12]}")
    return expected


def deck_base(html_path: Path) -> str:
    stem = html_path.stem
    return stem[:-6] if stem.endswith("_multi") else stem


def require_in_folder(path: Path, html_path: Path, folder_name: str) -> None:
    wanted = ("decks", folder_name, deck_base(html_path))
    legacy_wanted = ("html", "decks", folder_name, deck_base(html_path))
    parts = path.resolve().parts
    has_current = any(parts[i : i + 3] == wanted for i in range(len(parts) - 2))
    has_legacy = any(parts[i : i + 4] == legacy_wanted for i in range(len(parts) - 3))
    if not has_current and not has_legacy:
        fail(f"{path.name} must be stored under decks/{folder_name}/{deck_base(html_path)}/")


def validate_pdf(path: Path, html_path: Path, total: int) -> None:
    require_in_folder(path, html_path, "pdf")
    try:
        from pypdf import PdfReader
    except ImportError:
        fail("pypdf is not installed. Run `pip install -r .github/skills/ms-presentation-deck/scripts/requirements.txt`.")
    if not path.is_file() or path.stat().st_size == 0:
        fail(f"PDF missing or empty: {path}")
    reader = PdfReader(str(path))
    if len(reader.pages) != total:
        fail(f"PDF has {len(reader.pages)} pages, expected {total}")
    for index, page in enumerate(reader.pages, start=1):
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        if height == 0 or abs((width / height) - (16 / 9)) > 0.04:
            fail(f"PDF page {index} is not 16:9")
        if not (page.extract_text() or "").strip():
            fail(f"PDF page {index} has no extractable text")
    print(f"OK: PDF {path.name}, {len(reader.pages)} pages")


def note_matches(expected: str, actual: str) -> bool:
    expected_body = strip_markers(expected).lower()
    actual_body = strip_markers(actual).lower()
    return bool(expected_body and actual_body and expected_body[:80] in actual_body)


def validate_editable_shapes(slide, index: int, shape_type, full_width, full_height) -> None:
    shapes = list(slide.shapes)
    if not shapes:
        fail(f"PPTX slide {index} has no editable shapes")
    pictures = [shape for shape in shapes if shape.shape_type == shape_type.PICTURE]
    for picture in pictures:
        if picture.width >= full_width * 0.95 and picture.height >= full_height * 0.95:
            fail(f"PPTX slide {index} appears flattened as a full-slide image")
    editable_text = [shape for shape in shapes if getattr(shape, "has_text_frame", False) and normalize(shape.text)]
    if not editable_text:
        fail(f"PPTX slide {index} has no editable text")


def validate_slide_notes(slide, index: int, locale: str, expected: str) -> None:
    if not slide.has_notes_slide:
        fail(f"PPTX slide {index} has no notes slide")
    actual_note = normalize(slide.notes_slide.notes_text_frame.text)
    if not actual_note:
        fail(f"PPTX slide {index} has empty notes")
    if not any(marker in actual_note for marker in LOCALE_MARKERS[locale]):
        fail(f"PPTX slide {index} notes do not match locale {locale}")
    if not note_matches(expected, actual_note):
        fail(f"PPTX slide {index} notes do not match HTML source notes")


def validate_pptx(path: Path, html_path: Path, html: str, locale: str, total: int) -> None:
    require_in_folder(path, html_path, "pptx")
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        fail("python-pptx is not installed. Run `pip install -r .github/skills/ms-presentation-deck/scripts/requirements.txt`.")
    if not path.is_file() or path.stat().st_size == 0:
        fail(f"PPTX missing or empty: {path}")
    prs = Presentation(str(path))
    if len(prs.slides) != total:
        fail(f"PPTX has {len(prs.slides)} slides, expected {total}")
    expected = expected_notes(html, locale, total)
    full_width, full_height = prs.slide_width, prs.slide_height
    for index, slide in enumerate(prs.slides, start=1):
        validate_editable_shapes(slide, index, MSO_SHAPE_TYPE, full_width, full_height)
        validate_slide_notes(slide, index, locale, expected[index - 1])
    print(f"OK: PPTX {path.name}, {len(prs.slides)} editable slides with {locale} notes")


def main() -> int:
    parser = argparse.ArgumentParser(description="Validate deck PDF/PPTX derivatives")
    parser.add_argument("--html", required=True, help="Source HTML deck")
    parser.add_argument("--locale", required=True, choices=("en", "pt-BR", "es"))
    parser.add_argument("--pdf", help="PDF derivative path")
    parser.add_argument("--pptx", help="PPTX derivative path")
    parser.add_argument(
        "--allow-notes-derived-content",
        action="store_true",
        help="Allow PPTX visible slide text to be generated from localized speaker notes when the HTML visible text is not fully localized",
    )
    args = parser.parse_args()
    if not args.pdf and not args.pptx:
        fail("provide --pdf and/or --pptx")
    html_path = Path(args.html)
    if not html_path.is_file():
        fail(f"HTML source not found: {html_path}")
    html = html_path.read_text(encoding="utf-8")
    total = slide_count(html)
    if args.pdf and args.allow_notes_derived_content:
        fail("--allow-notes-derived-content is only allowed for PPTX, not PDF")
    if not args.allow_notes_derived_content:
        validate_visible_locale(html, args.locale)
    expected_notes(html, args.locale, total)
    if args.pdf:
        validate_pdf(Path(args.pdf), html_path, total)
    if args.pptx:
        validate_pptx(Path(args.pptx), html_path, html, args.locale, total)
    print("Derivative validation passed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
